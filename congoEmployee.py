# import pandas as pd
# from datetime import datetime
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# # Load Data
# df = pd.read_excel('employeeDetails.xlsx')
# df.columns = df.columns.str.strip()
# df['Joining Date'] = pd.to_datetime(df['Joining Date'])

# # Filter for today's anniversaries
# today = datetime.today()
# highlighted = df[
#     (df['Joining Date'].dt.month == today.month) &
#     (df['Joining Date'].dt.day == today.day)
# ]

# # Initialize presentation
# prs = Presentation()
# blank_slide_layout = prs.slide_layouts[6]

# # Constants for grid
# max_per_slide = 4
# image_width = Inches(2.5)
# image_height = Inches(2.3)
# name_height = Inches(0.4)
# cols = 2
# rows = 2
# padding_x = Inches(0.9)
# padding_y = Inches(1.2)
# spacing_x = Inches(6)
# spacing_y = Inches(3)

# # Function to apply background
# def apply_background(slide):
#     bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
#     bg.fill.solid()
#     bg.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light elegant blue
#     bg.line.fill.background()  # No border
#     slide.shapes._spTree.remove(bg._element)  # Send to back
#     slide.shapes._spTree.insert(2, bg._element)

# # Add title slide
# def add_title_slide():
#     slide = prs.slides.add_slide(blank_slide_layout)
#     apply_background(slide)

#     title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
#     tf = title_box.text_frame
#     tf.text = "🎉 Celebrating Our Employees"
#     p = tf.paragraphs[0]
#     p.font.size = Pt(44)
#     p.font.bold = True
#     p.font.color.rgb = RGBColor(0, 51, 102)

# # Add grid slide
# def add_grid_slide(employees):
#     slide = prs.slides.add_slide(blank_slide_layout)
#     apply_background(slide)

#     for i, row in enumerate(employees):
#         name = row['Employee Name']
#         photo = row['Photo']
#         col = i % cols
#         row_pos = i // cols

#         left = padding_x + col * spacing_x
#         top = padding_y + row_pos * spacing_y

#         # Add photo
#         try:
#             slide.shapes.add_picture(photo, left, top, image_width, image_height)
#         except:
#             # Fallback shape if image is missing
#             ph = slide.shapes.add_shape(
#                 MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, image_width, image_height
#             )
#             fill = ph.fill
#             fill.solid()
#             fill.fore_color.rgb = RGBColor(200, 200, 200)
#             tf = ph.text_frame
#             tf.text = "No Photo"
#             p = tf.paragraphs[0]
#             p.font.size = Pt(12)
#             p.font.color.rgb = RGBColor(80, 80, 80)

#         # Add name below the image
#         name_box = slide.shapes.add_textbox(left, top + image_height + Inches(0.1), image_width, name_height)
#         tf = name_box.text_frame
#         tf.text = name
#         p = tf.paragraphs[0]
#         p.font.size = Pt(12)
#         p.font.color.rgb = RGBColor(0, 0, 0)

# # Create slides
# add_title_slide()
# for i in range(0, len(highlighted), max_per_slide):
#     chunk = highlighted.iloc[i:i+max_per_slide].to_dict('records')
#     add_grid_slide(chunk)

# # Save presentation
# prs.save('Grid_Style_Employees.pptx')
# print("✅ Grid-style presentation saved as 'Grid_Style_Employees.pptx'")


import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from tkinter import Tk, Label, Button, filedialog
import os

# PPT generation logic (your existing code wrapped as a function)
def create_ppt(file_path):
    df = pd.read_excel(file_path)
    df.columns = df.columns.str.strip()
    df['Joining Date'] = pd.to_datetime(df['Joining Date'])

    today = datetime.today()
    highlighted = df[
        (df['Joining Date'].dt.month == today.month) &
        (df['Joining Date'].dt.day == today.day)
    ]

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    max_per_slide = 4
    image_width = Inches(2.5)
    image_height = Inches(2.3)
    name_height = Inches(0.4)
    cols = 2
    rows = 2
    padding_x = Inches(0.9)
    padding_y = Inches(1.2)
    spacing_x = Inches(6)
    spacing_y = Inches(3)

    def apply_background(slide):
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def add_title_slide():
        slide = prs.slides.add_slide(blank_slide_layout)
        apply_background(slide)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_box.text_frame
        tf.text = "🎉 Celebrating Our Employees"
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

    def add_grid_slide(employees):
        slide = prs.slides.add_slide(blank_slide_layout)
        apply_background(slide)

        for i, row in enumerate(employees):
            name = row['Employee Name']
            photo = row['Photo']
            col = i % cols
            row_pos = i // cols
            left = padding_x + col * spacing_x
            top = padding_y + row_pos * spacing_y

            try:
                slide.shapes.add_picture(photo, left, top, image_width, image_height)
            except:
                ph = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, image_width, image_height)
                fill = ph.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(200, 200, 200)
                tf = ph.text_frame
                tf.text = "No Photo"
                p = tf.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(80, 80, 80)

            name_box = slide.shapes.add_textbox(left, top + image_height + Inches(0.1), image_width, name_height)
            tf = name_box.text_frame
            tf.text = name
            p = tf.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0, 0, 0)

    add_title_slide()
    for i in range(0, len(highlighted), max_per_slide):
        chunk = highlighted.iloc[i:i+max_per_slide].to_dict('records')
        add_grid_slide(chunk)

    output_path = 'Grid_Style_Employees.pptx'
    prs.save(output_path)
    print("✅ Presentation created.")

    os.startfile(output_path)  # Automatically open the file (Windows only)

# GUI using Tkinter
def browse_file():
    file_path = filedialog.askopenfilename(
        title="Select Excel File",
        filetypes=[("Excel files", "*.xlsx *.xls")]
    )
    if file_path:
        create_ppt(file_path)
        status_label.config(text="✅ Presentation created and opened!")

# UI Setup
root = Tk()
root.title("Employee Anniversary PPT Generator")
root.geometry("400x200")

Label(root, text="Upload your Excel file", font=("Arial", 16)).pack(pady=20)
Button(root, text="Choose File", command=browse_file, width=20, bg="#0078D7", fg="white").pack(pady=10)
status_label = Label(root, text="", font=("Arial", 12), fg="green")
status_label.pack()

root.mainloop()
